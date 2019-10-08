import plotly.io as pio
import os
import cx_Oracle
import numpy as np
from datetime import date
from datetime import time
from datetime import datetime
from datetime import timedelta
import plotly.graph_objects as go
import pandas as pd
from IPython.display import Image
#from pd2ppt import df_to_powerpoint
from pptx import Presentation
from pptx.util import Inches
from pd2ppt import df_to_table
from pptx.util import Pt
from Guage_Graph import Create_Gauge
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
from datetime import date
import calendar
from pathlib import Path
# pio.orca.config.use_xvfb = False
pio.orca.config.use_xvfb = True
import plotly
plotly.io.orca.config.executable = '/usr/bin/orca'
plotly.io.orca.config.save()
pio.orca.config.save()


class DBQuery(object):

    def __init__(
            self,
            presentation_name
    ):

        self.query_list = []
        self.query_values = []
        self.path = str(Path(__file__).parent.absolute() + '/')
        self.presentation_name = presentation_name
        prs = Presentation(self.presentation_name)
        new_name = self.path + 'Reports/IBF_AUTOMATED_REPORT' + \
            str(datetime.now()) + '.pptx'
        prs.save(new_name)
        self.presentation_name = new_name
        
        
    def coverpage(self):
        t_left = Inches(8.5)
        t_top = Inches(1.3)
        t_width = Inches(4)
        t_height = Inches(0.7440945)
        prs = Presentation(self.presentation_name)
        slide = prs.slides[0]
        # Text
        txBox = slide.shapes.add_textbox(t_left, t_top, t_width, t_height)
        tf = txBox.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "IBF-MS-Daily Health "
        run.font.size = Pt(32)
        run.font.bold = True
        p = tf.add_paragraph()
        p.level = 0
        p.text = "Check Report:"
        p.font.size = Pt(32)
        p.font.bold = True
        p = tf.add_paragraph()
        p.level = 0
        today = date.today()
        now = datetime.now()
        dt_string = now.strftime("%d/%m/%Y %H:%M")
        print("date and time =", dt_string)
        p.text = dt_string
        p.font.size = Pt(32)
        p.font.bold = True

        prs.save(self.presentation_name)

    def ibs_kpi_020(self):

        self.query_list.append("IBS_KPI_020")
        dsn_tns = cx_Oracle.makedsn('10.206.11.82', '1521', sid='ibsdb1')
        conn = cx_Oracle.connect(
            user='mngd_svc', password='mngd_svc_201909', dsn=dsn_tns)

        cursor = conn.cursor()

        try:
            query_string = """
            select * from
            (
            select failed_svc, count(distinct(msisdn)) subsaffected
            from mtnibs_v2.bp_exceptions
            where trunc(svc_req_time) >= trunc(sysdate)-1 and trunc(svc_req_time) < trunc(sysdate)
            and severity not in ('RETRY', 'WARNING', '-')
            group by failed_svc
            order by subsaffected desc
            )
            """
            cursor.execute(query_string)

        except cx_Oracle.DatabaseError:
            conn.close()
            print('Failed to select from table')
            exit(1)

        query_result = cursor.fetchall()
        self.query_values.append(query_result)

        cursor.close()
        conn.close()

        query_result_df = pd.DataFrame(query_result)

        gauges = Create_Gauge(500)

        for index, row in query_result_df.iterrows():
            gauges.new_gauge(name=row[0], stat=row[1])

        gauges.create_full_image(name='ibs_kpi_020')

        # Put it in the presentation
        prs = Presentation(self.presentation_name)
        slide = prs.slides[7]
        shapes = prs.slides[0].shapes
        number_of_guages = gauges.count
        if gauges.count >= 10:
            picture = slide.shapes.add_picture(self.path + 'Images/ibs_kpi_020.jpg', left=Inches(0.7204724), top=Inches(1.775590551),
                                               width=Inches(6.732283), height=Inches(5.03937))
        elif gauges.columns < 10:
            picture = slide.shapes.add_picture(self.path + 'Images/ibs_kpi_020.jpg', left=Inches(1.2),
                                               top=Inches(1.775590551),
                                               width=Inches(5.15748), height=Inches(1.6*gauges.rows))

        prs.save(self.presentation_name)

        # This is the population of the explanations
        dsn_tns = cx_Oracle.makedsn('10.206.11.82', '1521', sid='ibsdb1')
        conn = cx_Oracle.connect(
            user='mngd_svc', password='mngd_svc_201909', dsn=dsn_tns)

        cursor = conn.cursor()
        count = 0
        message = ""
        for index, row in query_result_df.iterrows():
            top_margin = Inches(2.2 + (0.9 * count))
            try:
                query_string = """
                SELECT FAILED_ERR, DESCRIPTION, ROUND(COUNT( DISTINCT MSISDN)) averagefailures FROM
                (
                SELECT FAILED_SVC, FAILED_ERR, SEVERITY, DESCRIPTION, MSISDN
                FROM MTNIBS_V2.BP_EXCEPTIONS
                where trunc(svc_req_time) >= trunc(sysdate)-1 and trunc(svc_req_time) < trunc(sysdate)
                AND FAILED_SVC = '""" + row[0] + """'

                )
                GROUP BY FAILED_SVC, DESCRIPTION, FAILED_ERR
                ORDER BY averagefailures desc  
                        """
                cursor.execute(query_string)

            except cx_Oracle.DatabaseError:
                conn.close()
                print('Failed to select from table')
                exit(1)

            query_result = cursor.fetchall()
            query_result_df = pd.DataFrame(query_result)
            # @todo: pass this to a function to create the list and populate the presentation

            if count == 0 or count == 1:
                populate_list(slide, query_result_df, prs,
                              self.presentation_name, row[0], top_margin)

            print("Errors for :" + row[0])
            print(query_result_df)

            query_result_df.columns = [
                'ERROR No.', 'DESCRIPTION', 'AFFECTED SUBS.']
            message = message + "<br> Errors for: " + row[0] + "<br> "
            message = message + \
                query_result_df.to_html(index=False) + "<br>" + "<br>"
            count = count + 1
        self.query_values.append(query_result)

        cursor.close()
        conn.close()
        return message

    def ibs_kpi_021(self):

        self.query_list.append("IBS_KPI_021")
        dsn_tns = cx_Oracle.makedsn('10.206.11.82', '1521', sid='ibsdb1')
        conn = cx_Oracle.connect(
            user='mngd_svc', password='mngd_svc_201909', dsn=dsn_tns)

        cursor = conn.cursor()

        try:
            query_string = """
            SELECT * FROM
            (
            SELECT SEVERITY, ROUND(COUNT( DISTINCT MSISDN)) averagefailures
            FROM MTNIBS_V2.BP_EXCEPTIONS
            WHERE SVC_REQ_TIME > SYSDATE-1
            GROUP BY SEVERITY
            ORDER BY SEVERITY)
            WHERE averagefailures > 0
            """
            cursor.execute(query_string)

        except cx_Oracle.DatabaseError:
            conn.close()
            print('Failed to select from table')
            exit(1)

        query_result = cursor.fetchall()
        self.query_values.append(query_result)

        cursor.close()
        conn.close()

        query_result_df = pd.DataFrame(query_result)

        gauges = Create_Gauge(5000)
        # todo:need to sort the dataframe
        query_result_df = query_result_df.sort_values(1, ascending=False)
        for index, row in query_result_df.iterrows():
            gauges.new_gauge(name=row[0], stat=row[1])

        gauges.create_full_image(name='ibs_kpi_020')

        # Put it in the presentation
        prs = Presentation(self.presentation_name)
        slide = prs.slides[6]
        shapes = prs.slides[0].shapes
        number_of_guages = gauges.count

        picture = slide.shapes.add_picture(self.path + 'Images/ibs_kpi_020.jpg', left=Inches(1.2),
                                           top=Inches(1.775590551),
                                           width=Inches(5.15748), height=Inches(1.6*gauges.rows))

        prs.save(self.presentation_name)

        # This is the population of the explanations
        dsn_tns = cx_Oracle.makedsn('10.206.11.82', '1521', sid='ibsdb1')
        conn = cx_Oracle.connect(
            user='mngd_svc', password='mngd_svc_201909', dsn=dsn_tns)

        cursor = conn.cursor()
        count = 0
        message = ""
        for index, row in query_result_df.iterrows():
            # @todo: Create evaluation function to decide when i should put which rows
            top_margin = Inches(2.2 + (0.9 * count))
            try:
                query_string = """
                SELECT FAILED_ERR, DESCRIPTION, ROUND(COUNT( DISTINCT MSISDN)) averagefailures FROM
                (
                SELECT FAILED_ERR, SEVERITY, DESCRIPTION, MSISDN
                FROM MTNIBS_V2.BP_EXCEPTIONS
                WHERE SVC_REQ_TIME > SYSDATE-1
                AND SEVERITY = '""" + row[0] + """'

                )
                GROUP BY FAILED_ERR, DESCRIPTION
                ORDER BY averagefailures desc
                        """
                cursor.execute(query_string)

            except cx_Oracle.DatabaseError:
                conn.close()
                print('Failed to select from table')
                exit(1)

            query_result = cursor.fetchall()
            query_result_df = pd.DataFrame(query_result)
            # @todo: pass this to a function to create the list and populate the presentation

            if row[1] > 50:
                populate_list(slide, query_result_df, prs,
                              self.presentation_name, row[0], top_margin)

            print("Errors for: " + row[0])
            print(query_result_df)

            query_result_df.columns = [
                'ERROR No.', 'DESCRIPTION', 'AFFECTED SUBS.']
            message = message + "<br> Errors for: " + row[0] + "<br> "
            message = message + \
                query_result_df.to_html(index=False) + "<br>" + "<br>"
            count = count + 1
        self.query_values.append(query_result)

        cursor.close()
        conn.close()

        query_result_df = pd.DataFrame(query_result)
        return message

    def ibs_ms_001(self):

        self.query_list.append("IBS_MS_001")
        dsn_tns = cx_Oracle.makedsn('10.206.11.82', '1521', sid='ibsdb1')
        conn = cx_Oracle.connect(
            user='mngd_svc', password='mngd_svc_201909', dsn=dsn_tns)

        cursor = conn.cursor()
        my_date = date.today()
        dayOfWeek = datetime.today().strftime('%A')
        now = datetime.now()
        dateDiff = 1

        # If done on a monday then get it from friday
        # @todo: what to do if it's the end of day query?
        if now == 'Monday':
            dateDiff = 3
        begin_date_time = (now - timedelta(days=dateDiff)
                           ).strftime("%Y/%m/%d, %H:%M:%S")
        end_date_time = now.strftime("%Y/%m/%d, %H:%M:%S")
        try:
            query_string = """
            SELECT * FROM
            (
            select svcname, failed_svc, description, count(*) CNT1, count(distinct msisdn) affected_subscribers, CASE
                WHEN svcname <> ''
                THEN ''
                ELSE ''
            END COMMENTS
            from mtnibs_v2.bp_exceptions
            where (svc_req_time) between TO_DATE( '""" + str(begin_date_time) + """', 'YYYY/MM/DD HH24:MI:SS') and
            TO_DATE('""" + str(end_date_time) + """', 'YYYY/MM/DD HH24:MI:SS')
            and severity in ('FATAL','FAILED')
            group by svcname, failed_svc, description
            order by svcname, failed_svc, description
            )
            WHERE ROWNUM <= 1000 ORDER BY CNT1 DESC
            """
            cursor.execute(query_string)

        except cx_Oracle.DatabaseError:
            conn.close()
            print('Failed to select from table')
            exit(1)

        query_result = cursor.fetchall()
        self.query_values.append(query_result)

        cursor.close()
        conn.close()

        query_result_df = pd.DataFrame(query_result)
        query_result_df2 = query_result_df.transpose()
        query_result_df.columns = ['Service Name', 'Failed Component Service',
                                   'Error', 'Failure Count', 'Affected Subscribers', 'Comments']
        # fig = go.Figure(data=[go.Table(header=dict(values=['A Scores', 'B Scores', 'C', 'D', 'E', 'F']),
        #                                cells=dict(values=query_result_df2))])

        prs = Presentation(self.presentation_name)
        slide = prs.slides[5]
        now = datetime.now()

        dt_string = now.strftime("%d/%m/%Y %H:%M")
        text = 'Period: ' + str(dt_string) + ' - ' + str((now - timedelta(days=dateDiff)).strftime("%d/%m/%Y %H:%M"))
        txBox = slide.shapes.add_textbox(Inches(0.98), Inches(1.0), Inches(1.0), Inches(1.0))
        tf = txBox.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(12)
        top = Inches(2.5)
        left = Inches(0.75)
        width = Inches(9.25)
        height = Inches(5.0)
        shapes = slide.shapes

        tableObj = df_to_table(slide=slide, df=query_result_df, left=left,
                               top=top, width=width, name='ibs_kpi_001', height=1)
        table = tableObj.table

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(7)

        for cell in iter_header(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    run.font.bold = True

        table.columns[0].width = Inches(1.3)
        table.columns[1].width = Inches(1.6)
        table.columns[2].width = Inches(5.2)
        table.columns[3].width = Inches(0.8)
        table.columns[4].width = Inches(1.0)
        table.columns[5].width = Inches(1.5)

        prs.save(self.presentation_name)

    def ibs_ms_003(self):

        self.query_list.append("IBS_MS_003")
        dsn_tns = cx_Oracle.makedsn('10.211.11.78', '1521', sid='mcldb')
        conn = cx_Oracle.connect(
            user='ibs_sla', password='ibs_sla', dsn=dsn_tns)

        cursor = conn.cursor()

        try:
            query_string = """select a.*,
                                CASE
                                    WHEN severity <> ''
                                    THEN ''
                                    ELSE ''
                                END COMMENTS
                                from BP_EXCEPTION_STATS_VIEW a
                                where severity = 'RETRY'
                                order by severity
            """
            cursor.execute(query_string)

        except cx_Oracle.DatabaseError:
            conn.close()
            print('Failed to select from table ibs_ms_003')
            exit(1)

        query_result = cursor.fetchall()
        self.query_values.append(query_result)

        cursor.close()
        conn.close()

        query_result_df = pd.DataFrame(query_result)
        services = query_result_df.iloc[:, 1]

        services = query_result_df.iloc[:, 1]
        x = query_result_df.iloc[:, 2]
        xname = 'Unique MSISDNs Affected (2 Days Ago)'
        y = query_result_df.iloc[:, 3]
        yname = 'Unique MSISDNs Affected (Yesterday)'
        z = query_result_df.iloc[:, 4]
        zname = 'Unique MSISDNs Affected (Today)'

        df = pd.DataFrame({xname: x.as_matrix(), yname: y.as_matrix(
        ), zname: z.as_matrix()}, index=services)
        ax = df.plot.barh(figsize=(10, 6))
        fig = ax.get_figure()
        fig.savefig(self.path + 'Images/ibs_ms_003.png', transparent=True)

        prs = Presentation(self.presentation_name)
        slide = prs.slides[3]
        top = Inches(2.0)
        left = Inches(0.25)
        width = Inches(9.25)
        height = Inches(5.0)
        shapes = slide.shapes
        picture = slide.shapes.add_picture(self.path + 'Images/ibs_ms_003.png', left=Inches(0.9),
                                           top=Inches(1.3),
                                           width=Inches(10.23), height=Inches(6.02))
        prs.save(self.presentation_name)

        slide = prs.slides[4]
        query_result_df.columns = ['Severity', 'Service Name', '2 Days Ago', '1 Day Ago', 'Today',
                                   'Total Invocations 2 Days Ago', 'Total Invocations 1 Days Ago',
                                   'Total Invocations Today', 'Retry % 2 Days Ago', 'Retry % 1 Day Ago',
                                   'Retry % Today', 'Comments']

        tableObj = df_to_table(slide=slide, df=query_result_df,
                               left=left, top=top, width=width, name='ibs_kpi_003')
        table = tableObj.table

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(6)

        for cell in iter_header(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    run.font.bold = True

        table.columns[0].width = Inches(1.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(1.0)
        table.columns[3].width = Inches(1.0)
        table.columns[4].width = Inches(1.0)
        table.columns[5].width = Inches(1.0)
        table.columns[6].width = Inches(1.0)
        table.columns[7].width = Inches(1.0)
        table.columns[8].width = Inches(1.0)
        table.columns[9].width = Inches(1.0)
        table.columns[10].width = Inches(1.0)
        table.columns[11].width = Inches(0.9)
        prs.save(self.presentation_name)


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def iter_header(table):
    row = table.rows[0]
    for cell in row.cells:
        yield cell


def populate_list(slide, query_result_df, prs, name, severity, top):
    # Text position
    t_left = Inches(8.5)
    t_top = top
    t_width = Inches(4)
    t_height = Inches(0.7440945)
    # Text
    txBox = slide.shapes.add_textbox(t_left, t_top, t_width, t_height)
    tf = txBox.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = severity
    run.font.size = Pt(10)
    p = tf.add_paragraph()
    p.level = 0
    p.text = "\u2022  " + " These are based on the following exception/error:"
    p.font.size = Pt(9)
    p = tf.add_paragraph()
    p.level = 0
    p.text = "     " + query_result_df[0][0] + ": " + query_result_df[1][0]
    p.font.size = Pt(9)
    p = tf.add_paragraph()
    p.level = 0
    data = pd.read_csv("Presentation/Errors.csv", sep='|')

    try:
        err_description = data[data.ERROR == str(query_result_df[0][0])].iloc[0][1]
    except:
        err_description = '-'
    p.text = "     Error " + query_result_df[0][0] + ": " + str(err_description)
    p.font.size = Pt(9)
    prs.save(name)
